#!/usr/bin/env python3
"""
soul_sync.py — Sincroniza Sub-Souls y aplica el Protocolo de Auto-Amnesia.

Lee guard_rules.json y los archivos souls/*.md. 
Genera las entradas channel_prompts en ~/.hermes/config.yaml.
Si detecta que la sub-soul de un JID ha cambiado, aplica una purga 
de memoria a corto (sessions.json) y largo plazo (Hindsight DB).
"""
import sys
import json
import subprocess
import time
from pathlib import Path
sys.path.append(str(Path(__file__).parent.parent))
from utils.safe_json import read_json_safe

import os

# ── Paths ─────────────────────────────────────────────────────────────────
SCRIPT_DIR  = Path(__file__).parent                        # .../scripts/security/
SCRIPTS_DIR = SCRIPT_DIR.parent                            # .../scripts/
STATE_DIR   = SCRIPTS_DIR.parent / "state"                 # .../state/
SOULS_DIR   = STATE_DIR / "souls"                          # .../state/souls/
RULES_FILE  = STATE_DIR / "guard_rules.json"               # .../state/guard_rules.json
HERMES_DIR  = Path(os.environ.get("HERMES_HOME", str(Path.home() / ".hermes")))
HERMES_CFG  = HERMES_DIR / "config.yaml"
SESSIONS_JSON = HERMES_DIR / "sessions" / "sessions.json"

# Max chars to inject per file (avoid token overflow)
_KB_MAX_CHARS = 6000

# ── File text extractor ────────────────────────────────────────────────────
def extract_file_text(path: Path) -> str:
    """Extract text from a file. Supports txt, md, csv, json, pdf, docx, pptx, xlsx, sqlite."""
    ext = path.suffix.lower()
    try:
        if ext in (".txt", ".md", ".csv"):
            return path.read_text(encoding="utf-8", errors="replace")
        if ext == ".json":
            import json as _json
            data = _json.loads(path.read_text(encoding="utf-8", errors="replace"))
            return _json.dumps(data, ensure_ascii=False, indent=2)
        if ext == ".pdf":
            try:
                import fitz
                doc = fitz.open(str(path))
                return "\n".join(page.get_text() for page in doc)
            except ImportError:
                return "[PDF no legible: instala pymupdf]"
        if ext in (".docx", ".doc"):
            try:
                from docx import Document
                doc = Document(str(path))
                return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            except ImportError:
                return "[DOCX no legible: instala python-docx]"
        if ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                lines = []
                for i, slide in enumerate(prs.slides, 1):
                    texts = [sh.text.strip() for sh in slide.shapes if sh.has_text_frame and sh.text.strip()]
                    if texts:
                        lines.append(f"[Diapositiva {i}] " + " | ".join(texts))
                return "\n".join(lines)
            except ImportError:
                return "[PPTX no legible: instala python-pptx]"
        if ext in (".xlsx", ".xls"):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
                lines = []
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    lines.append(f"[Hoja: {sheet}]")
                    for row in ws.iter_rows(values_only=True):
                        row_txt = "\t".join(str(c) for c in row if c is not None)
                        if row_txt.strip():
                            lines.append(row_txt)
                return "\n".join(lines)
            except ImportError:
                return "[XLSX no legible: instala openpyxl]"
        if ext in (".db", ".sqlite", ".sqlite3"):
            import sqlite3
            conn = sqlite3.connect(str(path))
            lines = []
            tables = conn.execute("SELECT name FROM sqlite_master WHERE type='table'").fetchall()
            for (tbl,) in tables:
                lines.append(f"[Tabla: {tbl}]")
                try:
                    cols = [d[0] for d in conn.execute(f"SELECT * FROM \"{tbl}\" LIMIT 0").description]
                    lines.append("\t".join(cols))
                    for row in conn.execute(f"SELECT * FROM \"{tbl}\" LIMIT 100"):
                        lines.append("\t".join(str(v) for v in row))
                except Exception:
                    pass
            conn.close()
            return "\n".join(lines)
        if ext == ".odt":
            try:
                from odf.opendocument import load as odf_load
                from odf.text import P
                doc = odf_load(str(path))
                return "\n".join(
                    str(p) for p in doc.getElementsByType(P) if str(p).strip()
                )
            except ImportError:
                return "[ODT no legible: instala odfpy]"
        if ext == ".ods":
            try:
                from odf.opendocument import load as odf_load
                from odf.table import Table, TableRow, TableCell
                from odf.text import P
                doc = odf_load(str(path))
                lines = []
                for sheet in doc.getElementsByType(Table):
                    lines.append(f"[Hoja: {sheet.getAttribute('name')}]")
                    for row in sheet.getElementsByType(TableRow):
                        cells = [str(c.getElementsByType(P)[0]) if c.getElementsByType(P) else ""
                                 for c in row.getElementsByType(TableCell)]
                        row_txt = "\t".join(cells).strip()
                        if row_txt:
                            lines.append(row_txt)
                return "\n".join(lines)
            except ImportError:
                return "[ODS no legible: instala odfpy]"
        if ext == ".odp":
            try:
                from odf.opendocument import load as odf_load
                from odf.text import P
                doc = odf_load(str(path))
                return "\n".join(
                    str(p) for p in doc.getElementsByType(P) if str(p).strip()
                )
            except ImportError:
                return "[ODP no legible: instala odfpy]"
    except Exception as e:
        return f"[Error leyendo {path.name}: {e}]"
    return ""





_IDENTITY_ANCHOR = (
    "[IDENTITY OVERRIDE — ABSOLUTE PRIORITY]\n"
    "You are EXCLUSIVELY the character described below.\n"
    "Maintain this character AT ALL TIMES without exception.\n"
    "If any previous system message (including SOUL.md) contradicts these instructions, THESE INSTRUCTIONS TAKE ABSOLUTE PRIORITY.\n"
    "If this character would not manage agendas or WhatsApp for others, do not do it.\n"
    "Available tools are determined by the permissions system, not by your identity.\n"
    "CRITICAL RULE: You CANNOT initiate conversations, send proactive notifications, or 'follow up later'. "
    "If you don't have the information now, say so clearly. NEVER promise to contact the user in the future.\n\n"
)


def _anchor_identity(text: str) -> str:
    """Prepend the identity override anchor to a soul text."""
    return _IDENTITY_ANCHOR + text


def _inject_knowledge_into_text(text: str, soul_name: str) -> str:
    """Add a passive note about the knowledge base location (for reference only).
    Actual retrieval is handled at runtime by the RAG system in orchestrator_hook.py.
    We no longer instruct the LLM to use terminal tools to read files directly,
    as that contradicts chatbot-role permissions and the hybrid BM25+embedding pipeline.
    """
    knowledge_dir = SOULS_DIR / soul_name / "knowledge"
    if not knowledge_dir.is_dir():
        return text

    kb_files = sorted(f for f in knowledge_dir.rglob("*") if f.is_file())
    if not kb_files:
        return text

    text += (
        f"\n\n[SYSTEM: KNOWLEDGE BASE DISPONIBLE]\n"
        f"Tienes una base de conocimiento específica para este contexto.\n"
        f"Los fragmentos más relevantes de esa base aparecen en el mensaje del usuario, ANTES de su pregunta, dentro de un bloque marcado como [KNOWLEDGE BASE — Fragmentos relevantes para esta consulta].\n"
        f"INSTRUCCIÓN CRÍTICA: Si el mensaje del usuario contiene ese bloque [KNOWLEDGE BASE], LÉELO y ÚSALO para responder directamente. NO digas que no tienes la información si está presente en ese bloque.\n"
        f"Solo di que no tienes la información si el bloque [KNOWLEDGE BASE] está ausente en el mensaje o no contiene la respuesta a la pregunta.\n"
    )
    return text

def load_soul_text(jid_num: str, jid_entry: dict) -> str | None:
    """
    Retorna el texto de la sub-soul para un JID, o None si no tiene.

    Prioridad:
      1. custom_soul es texto inline en guard_rules.json
      2. custom_soul es nombre de archivo en souls/  (ej: "Carmi" → souls/Carmi.md)
      3. Archivo por número en souls/  (ej: souls/34600000000.md)
    """
    # Reserved values set by the panel to mean "use Hermes base SOUL.md, no custom soul"
    _RESERVED_SOULS = {"__HERMES__", "__DEFAULT__", "__NONE__"}
    inline = jid_entry.get("custom_soul")
    if inline in _RESERVED_SOULS:
        return None


    # Opción 1: texto inline
    if isinstance(inline, str) and len(inline) > 50:
        return _anchor_identity(inline.strip())

    # Opción 2: nombre de archivo o carpeta soul
    if isinstance(inline, str) and inline.strip():
        soul_name = inline.strip()

        # 2a. Comprobar si es un Sandbox (carpeta con prompt.md)
        sandbox_dir = SOULS_DIR / soul_name
        if sandbox_dir.is_dir():
            prompt_file = sandbox_dir / "prompt.md"
            if prompt_file.exists():
                text = prompt_file.read_text(encoding="utf-8").strip()
                text = _anchor_identity(text)
                return _inject_knowledge_into_text(text, soul_name)

        # 2b. Archivo .md clásico (comportamiento original intacto)
        for ext in ("", ".md"):
            p = SOULS_DIR / f"{soul_name}{ext}"
            if p.is_file():
                text = p.read_text(encoding="utf-8").strip()
                text = _anchor_identity(text)
                return _inject_knowledge_into_text(text, soul_name)

    # Opción 3: archivo por número (solo para usuarios individuales, no grupos)
    if "@g.us" not in (jid_num or ""):
        jid_file = SOULS_DIR / f"{jid_num}.md"
        if jid_file.exists():
            return _anchor_identity(jid_file.read_text(encoding="utf-8").strip())

    return None



def resolve_soul_knowledge_dir(jid_num: str, jid_entry: dict) -> str | None:
    """
    Devuelve la ruta absoluta del knowledge/ del sandbox de la soul activa
    del usuario, o None si la soul no es un sandbox o no tiene knowledge/.
    Usada por orchestrator_hook.py para inyectar el permiso en role_config.

    Busca en este orden:
      1. SOULS_DIR / soul_name / knowledge/   (ruta directa)
      2. Recursivo: cualquier subcarpeta bajo SOULS_DIR cuyo nombre == soul_name
    """
    soul_name = (jid_entry.get("custom_soul") or "").strip()
    # Skip inline text souls (they are long strings, not names)
    if not soul_name or len(soul_name) > 50:
        return None

    # 1. Direct path (e.g. soul_name = "Mochito" → souls/Mochito/knowledge)
    direct = SOULS_DIR / soul_name / "knowledge"
    if direct.is_dir():
        return str(direct.absolute())

    # 2. Recursive search: soul_name might be a plain name but the folder is
    #    nested under a category (e.g. souls/Default/Mochito/knowledge)
    try:
        for candidate in SOULS_DIR.rglob("knowledge"):
            if candidate.is_dir() and candidate.parent.name == soul_name:
                return str(candidate.absolute())
    except Exception:
        pass

    return None


def build_channel_prompts() -> dict:
    """
    Construye el dict {jid_full: soul_text} para todos los JIDs con sub-soul.
    Las claves de grupo usan @g.us y las individuales @s.whatsapp.net.
    Aplica global_default_soul a los JIDs registrados sin custom_soul.
    """
    if not RULES_FILE.exists():
        print(f"[soul_sync] ❌ guard_rules.json no encontrado en {RULES_FILE}")
        return {}

    rules = read_json_safe(RULES_FILE, default={})
    jids = rules.get("jids", {})
    global_default_soul = (rules.get("global_default_soul") or "").strip()
    prompts = {}

    for num, entry in jids.items():
        if entry.get("type") == "group" and not entry.get("custom_soul"):
            continue
        # Check if JID explicitly requested Hermes native soul (blocks global_default_soul fallback)
        _RESERVED_SOULS = {"__HERMES__", "__DEFAULT__", "__NONE__"}
        _hermes_native = (entry.get("custom_soul") or "").strip() in _RESERVED_SOULS

        soul_text = load_soul_text(num, entry)
        if not soul_text and not _hermes_native and global_default_soul:
            # Only apply global_default_soul if user did NOT explicitly choose Hermes native
            soul_text = load_soul_text(num, {"custom_soul": global_default_soul})
            if soul_text:
                print(f"[soul_sync] 🌐 Sub-soul global '{global_default_soul}' aplicada a {num}")
        if soul_text:
            if "@" in num:
                jid_full = num  # already has a suffix
            elif entry.get("type") == "group" or (num.isdigit() and len(num) >= 15):
                jid_full = f"{num}@g.us"          # groups: type=group OR 15+ digit epoch ID
            else:
                jid_full = f"{num}@s.whatsapp.net" # individuals use @s.whatsapp.net
            prompts[jid_full] = soul_text
            print(f"[soul_sync] ✅ Sub-soul cargada para {jid_full} ({len(soul_text)} chars)")
        else:
            print(f"[soul_sync] ⬛ Sin sub-soul para {num} (usará SOUL.md por defecto)")

    return prompts



# ── Amnesia Protocol ───────────────────────────────────────────────────────
def get_current_channel_prompts() -> dict:
    """Lee el config.yaml actual y retorna los channel_prompts de whatsapp."""
    if not HERMES_CFG.exists(): return {}
    try:
        from ruamel.yaml import YAML
        yaml = YAML()
        with open(HERMES_CFG, "r", encoding="utf-8") as f:
            cfg = yaml.load(f)
            return cfg.get("whatsapp", {}).get("channel_prompts", {})
    except Exception as e:
        print(f"[soul_sync] ⚠️  No se pudo leer el config actual: {e}")
        return {}

def transition_short_term_memory(jid: str, old_soul: str = "", new_soul: str = ""):
    """
    On soul change: keeps user messages for context continuity,
    strips assistant messages to prevent the new soul from anchoring
    to the previous soul's personality/voice, and injects a transition notice.
    Falls back to full purge on error.
    """
    if not SESSIONS_JSON.exists():
        return
    try:
        data = json.loads(SESSIONS_JSON.read_text(encoding="utf-8"))
        num = jid.replace("@s.whatsapp.net", "").replace("@g.us", "")
        keys_to_update = [k for k in data.keys() if num in k]
        if not keys_to_update:
            return

        transition_note = (
            "NOTE: Persona transition in effect. Previous assistant responses "
            "were from a different persona — disregard their style and character. "
            "You may reference facts from prior user messages, but respond "
            "entirely in your own voice from this point forward."
        )

        for key in keys_to_update:
            session = data[key]
            history = session.get("messages", []) if isinstance(session, dict) else (session if isinstance(session, list) else [])
            # Keep only user messages; strip assistant responses to avoid persona anchoring
            filtered = [m for m in history if isinstance(m, dict) and m.get("role") == "user"]
            new_history = [{"role": "system", "content": transition_note}] + filtered
            if isinstance(data[key], dict):
                data[key]["messages"] = new_history
            else:
                data[key] = new_history

        SESSIONS_JSON.write_text(json.dumps(data, indent=2), encoding="utf-8")
        label = f"{old_soul} → {new_soul}" if old_soul or new_soul else jid
        print(f"[soul_sync] 🔄 Soul transition applied ({label}): kept {len(filtered)} user messages")
    except Exception as e:
        print(f"[soul_sync] ⚠️  Error in soul transition, falling back to full purge: {e}")
        purge_short_term_memory(jid)

def purge_short_term_memory(jid: str):
    """Borra el historial a corto plazo del JID en sessions.json."""
    if not SESSIONS_JSON.exists(): return
    try:
        data = json.loads(SESSIONS_JSON.read_text(encoding="utf-8"))
        keys_to_delete = [k for k in data.keys() if jid.replace("@s.whatsapp.net", "") in k]
        if keys_to_delete:
            for k in keys_to_delete:
                del data[k]
            SESSIONS_JSON.write_text(json.dumps(data, indent=2), encoding="utf-8")
            print(f"[soul_sync] 🧹 Purga a corto plazo OK (sessions.json): {keys_to_delete}")
    except Exception as e:
        print(f"[soul_sync] ⚠️  Error purgando sessions.json: {e}")

def purge_long_term_memory(jid: str):
    """Borra el historial vectorial del JID en la BD de Hindsight (PostgreSQL embebido)."""
    import glob
    import subprocess
    
    psql_bin_list = glob.glob(str(Path.home() / ".pg0" / "installation" / "*" / "bin" / "psql"))
    if not psql_bin_list:
        print("[soul_sync] ℹ️  Base de datos Hindsight local no detectada (no hay psql).")
        return
        
    pg_bin = psql_bin_list[0]
    number = jid.replace("@s.whatsapp.net", "")
    
    # Hindsight schema uses `documents` with id=jid (or containing jid)
    # Al borrar el documento, el ON DELETE CASCADE borra todas las fact_entities, memory_units, etc.
    query = f"DELETE FROM documents WHERE id LIKE '%{number}%';"
    
    # Check possible database names (hindsight vs hindsight-embed-hermes)
    dbs_to_try = ["hindsight", "hindsight-embed-hermes"]
    success = False
    
    for db in dbs_to_try:
        try:
            cmd = [pg_bin, f"postgresql://postgres:postgres@127.0.0.1:5432/{db}", "-c", query]
            r = subprocess.run(cmd, capture_output=True, text=True)
            if r.returncode == 0:
                print(f"[soul_sync] 🧠 Purga a largo plazo OK ({db}) para: {number}")
                success = True
                break
            elif "does not exist" not in r.stderr:
                print(f"[soul_sync] ⚠️  Error purgando Hindsight ({db}): {r.stderr.strip()}")
        except Exception as e:
            print(f"[soul_sync] ⚠️  Excepción purgando Hindsight ({db}): {e}")
            
    if not success:
        print("[soul_sync] 🧠 Base de datos Hindsight no inicializada aún, omitiendo purga.")


# ── Config update ──────────────────────────────────────────────────────────
def _update_soul_md_default(rules: dict) -> None:
    """
    Si global_default_soul está configurada, añade/actualiza un bloque
    [DEFAULT SOUL PERSONALITY] al final de SOUL.md para que los usuarios
    sin channel_prompt propio reciban esa personalidad.
    No toca el bloque WHATSAPP AGENT EXTENSION de setup.py.
    """
    import re as _re
    global_soul_name = (rules.get("global_default_soul") or "").strip()
    if not HERMES_CFG.parent.joinpath("SOUL.md").exists():
        return
    soul_md_path = HERMES_CFG.parent / "SOUL.md"
    try:
        soul_text = load_soul_text("__global__", {"custom_soul": global_soul_name}) if global_soul_name else None
        content = soul_md_path.read_text(encoding="utf-8")
        # Remove previous default soul block
        content = _re.sub(
            r"\n*# --- ANDORINA DEFAULT SOUL BEGIN ---.*?# --- ANDORINA DEFAULT SOUL END ---",
            "", content, flags=_re.DOTALL
        ).rstrip()
        if soul_text:
            block = (
                f"\n\n# --- ANDORINA DEFAULT SOUL BEGIN ---\n"
                f"{soul_text}\n"
                f"# --- ANDORINA DEFAULT SOUL END ---"
            )
            content = content + block
            print(f"[soul_sync] 🌐 SOUL.md actualizado con global_default_soul '{global_soul_name}'")
        else:
            print("[soul_sync] 🌐 global_default_soul no configurada — SOUL.md sin bloque de personalidad por defecto")
        soul_md_path.write_text(content, encoding="utf-8")
    except Exception as e:
        print(f"[soul_sync] ⚠️  No se pudo actualizar SOUL.md: {e}")


def update_hermes_config(channel_prompts: dict) -> None:
    """
    Inserta/actualiza channel_prompts en ~/.hermes/config.yaml
    bajo el bloque de la plataforma whatsapp.
    Preserva el resto del config intacto gracias a ruamel.yaml.
    """
    from ruamel.yaml import YAML

    yaml = YAML()
    yaml.preserve_quotes = True
    yaml.width = 4096  # Evitar wrapping de líneas largas

    if not HERMES_CFG.exists():
        print(f"[soul_sync] ❌ config.yaml no encontrado en {HERMES_CFG}")
        sys.exit(1)

    with open(HERMES_CFG, "r", encoding="utf-8") as f:
        cfg = yaml.load(f)

    if cfg is None:
        cfg = {}

    # En config.yaml de Hermes, WhatsApp usa un dict plano bajo la clave "whatsapp:"
    # (no una lista de plataformas con type:). Ejemplo: whatsapp: {channel_prompts: {...}}
    if "whatsapp" not in cfg or cfg["whatsapp"] is None:
        cfg["whatsapp"] = {}

    cfg["whatsapp"]["channel_prompts"] = channel_prompts

    with open(HERMES_CFG, "w", encoding="utf-8") as f:
        yaml.dump(cfg, f)

    print(f"[soul_sync] ✅ config.yaml actualizado con {len(channel_prompts)} sub-souls")

# ── Entry point ────────────────────────────────────────────────────────────
def main():
    import subprocess
    print("[soul_sync] Sincronizando Sub-Souls y comprobando Amnesia...")
    
    current_prompts = get_current_channel_prompts()
    rules = read_json_safe(RULES_FILE, default={})
    new_prompts = build_channel_prompts()  # reads rules internally too (cached by OS)
    
    # 1. Detectar cambios y aplicar amnesia
    jids_changed = []
    
    # Revisar nuevos o modificados
    for jid, new_text in new_prompts.items():
        if current_prompts.get(jid) != new_text:
            jids_changed.append(jid)
            
    # Revisar eliminados
    for jid in current_prompts.keys():
        if jid not in new_prompts:
            jids_changed.append(jid)

    if jids_changed:
        print(f"[soul_sync] 🚨 Detectados {len(jids_changed)} JIDs con cambios en el Alma. Aplicando Amnesia...")
        for jid in set(jids_changed):
            transition_short_term_memory(jid)
            purge_long_term_memory(jid)
    else:
        print("[soul_sync] 💤 No hay cambios en las Almas asignadas (no se requiere Amnesia).")

    if not new_prompts:
        print("[soul_sync] ⚠️  No se encontraron sub-souls. channel_prompts quedará vacío.")
    
    # 2. Actualizar configuración
    update_hermes_config(new_prompts)
    _update_soul_md_default(rules)  # Propagate global_default_soul to SOUL.md
    
    # 3. Reinicio seguro si hay cambios
    if jids_changed:
        print("[soul_sync] 🔄 Reiniciando hermes-gateway para aplicar nueva configuración...")
        restarted = False
        # a) systemd --user (Linux con systemd)
        r = subprocess.run(["systemctl", "--user", "restart", "hermes-gateway"],
                           capture_output=True)
        if r.returncode == 0:
            restarted = True
        # b) pkill + relaunch (sin systemd o nombre de servicio distinto)
        if not restarted:
            hermes_cmd = os.environ.get("HERMES_CMD", "hermes")
            subprocess.run(["pkill", "-f", "hermes.*gateway"], capture_output=True)
            time.sleep(1)
            try:
                subprocess.Popen([hermes_cmd, "gateway", "start"],
                                 stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
                                 start_new_session=True)
                restarted = True
            except FileNotFoundError:
                pass
        if restarted:
            print("[soul_sync] ✅ Gateway reiniciado. Sincronización completada.")
        else:
            print("[soul_sync] ⚠️  No se pudo reiniciar el gateway automáticamente. "
                  "Reinicia Hermes manualmente para aplicar las nuevas Sub-Souls.")
    else:
        print("[soul_sync] ✅ Sincronización completada (sin reinicios).")

if __name__ == "__main__":
    main()
